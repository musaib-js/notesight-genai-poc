from fastapi import FastAPI, UploadFile, File, Form, APIRouter
from fastapi.responses import StreamingResponse
import os
import asyncio
import logging
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation
import pandas as pd
from openai import AsyncOpenAI
from mistralai import Mistral
from core.config import OPENAI_API_KEY, MISTRAL_API_KEY, GEMINI_API_KEY
from core.prompts import SUMMARY_PROMPT
from google import genai
import base64

router = APIRouter()
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

openai_client = AsyncOpenAI(api_key=OPENAI_API_KEY)
mistral_client = Mistral(api_key=MISTRAL_API_KEY)

async def extract_text(pdf_path: str, start_page: int, end_page: int) -> str:
    """Extract text from PDF pages."""
    with pdfplumber.open(pdf_path) as pdf:
        text = "\n".join(pdf.pages[i].extract_text() or "" for i in range(start_page, min(end_page, len(pdf.pages))))
    return text

async def generate_notes_stream_chatgpt(cleaned_text: str, previous_summary: str = ""):
    """Generate structured notes using OpenAI with streaming."""
    if not cleaned_text:
        yield ""

    try:
        messages = [{"role": "user", "content": SUMMARY_PROMPT.format(text=cleaned_text)}]

        if previous_summary:
            messages.insert(0, {"role": "assistant", "content": previous_summary})

        response = await openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            max_tokens=1000,
            temperature=0.7,
            stream=True
        )

        async for chunk in response:
            if chunk.choices[0].delta.content:
                yield chunk.choices[0].delta.content

    except Exception as e:
        logging.error(f"ChatGPT Streaming Error: {e}")
        yield f"Error: {str(e)}"

async def generate_notes_stream_mistral(cleaned_text: str,previous_summary: str = ""):
    """Generate structured notes using Mistral AI with streaming."""
    try:
        messages = [{"role": "user", "content": SUMMARY_PROMPT.format(text=cleaned_text)}]

        if previous_summary:
            messages.insert(0, {"role": "assistant", "content": previous_summary})
        response = await mistral_client.chat.stream_async(
            model="mistral-medium",
            messages=messages
        )

        async for chunk in response:
            if chunk.data.choices[0].delta.content:
                yield chunk.data.choices[0].delta.content

    except Exception as e:
        logging.error(f"Mistral Streaming Error: {e}")
        yield f"Error: {str(e)}"

async def generate_gemini_notes_stream(cleaned_text: str, previous_summary: str = ""):
    """Generate structured notes using Gemini AI with streaming."""
    client = genai.Client(api_key=GEMINI_API_KEY, http_options={'api_version': 'v1alpha'})
    model = "gemini-2.0-flash-exp"
    config = {"response_modalities": ["TEXT"]}
    if not cleaned_text:
        yield ""

    try:
        messages = SUMMARY_PROMPT.format(text=cleaned_text)

        if previous_summary:
            messages = previous_summary + "\n\n" + messages
        async with client.aio.live.connect(model=model, config=config) as session:
            message = messages

            await session.send(input=message, end_of_turn=True)

            full_summary = ""

            async for response in session.receive():
                if response.text is not None:
                    full_summary += response.text
                    yield response.text

    except Exception as e:
        logging.error(f"Gemini Streaming Error: {e}")
        yield f"Error: {str(e)}"

async def process_chunk(pdf_path: str, start_page: int, end_page: int):
    """Extract and clean text from a chunk of pages asynchronously."""
    raw_text = await extract_text(pdf_path, start_page, end_page)
    cleaned_text = raw_text.replace("-\n", "").replace("\n", " ").strip()
    return cleaned_text if cleaned_text else None

def encode_image(image_path):
    """Encode the image to base64."""
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except FileNotFoundError:
        print(f"Error: The file {image_path} was not found.")
        return None
    except Exception as e:  
        print(f"Error: {e}")
        return None

def extract_text_from_file(file_path: str) -> str:
    """Extract text from different file types."""
    ext = os.path.splitext(file_path)[-1].lower()
    text = ""
    
    if ext == ".pdf":
        with pdfplumber.open(file_path) as pdf:
            text = "\n".join([page.extract_text() or "" for page in pdf.pages])
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    elif ext in [".jpg", ".png"]:
        mistral_client = Mistral(api_key=MISTRAL_API_KEY)
        base64_image = encode_image(file_path)
        ocr_response =mistral_client.ocr.process(
            model="mistral-ocr-latest",
            document={"type": "image_url", "image_url":f"data:image/jpeg;base64,{base64_image}"}
        )
        text = "\n".join(page.markdown for page in ocr_response.pages)
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = "\n".join([shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
    elif ext in [".xlsx", ".csv"]:
        df = pd.read_excel(file_path, dtype=str) if ext == ".xlsx" else pd.read_csv(file_path, dtype=str)
        text = df.to_string(index=False)  
    
    return text.strip()

def chunk_text(text: str, chunk_size: int = 10000):
    """Split text into smaller chunks."""
    words = text.split()
    for i in range(0, len(words), chunk_size):
        yield " ".join(words[i:i + chunk_size])

async def stream_summary(file_paths: list[str], model: str):
    """Stream summarized notes for multiple files using OpenAI, Mistral, or Gemini."""
    previous_summary = ""
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[-1].lower()
        if ext == ".pdf":
            try:
                with pdfplumber.open(file_path) as pdf:
                    total_pages = len(pdf.pages)

                for start in range(0, total_pages, 5):
                    cleaned_text = await process_chunk(file_path, start, min(start + 5, total_pages))

                    if cleaned_text:
                        if model == "chatgpt":
                            async for chunk in generate_notes_stream_chatgpt(cleaned_text, previous_summary):
                                yield chunk
                            previous_summary = chunk

                        elif model == "mistral":
                            async for chunk in generate_notes_stream_mistral(cleaned_text):
                                yield chunk
                            previous_summary = chunk

                        elif model == "gemini":
                            async for chunk in generate_gemini_notes_stream(cleaned_text,previous_summary):
                                yield chunk
                            previous_summary = chunk
            except Exception as e:
                logging.error(f"Streaming error: {e}")
                yield f"Error: {str(e)}"
        else:
            try:
                cleaned_text = extract_text_from_file(file_path)
                if cleaned_text:
                    for chunk in chunk_text(cleaned_text):
                        if model == "chatgpt":
                            async for response in generate_notes_stream_chatgpt(chunk, previous_summary):
                                yield response
                            previous_summary = response
                        elif model == "mistral":
                            async for response in generate_notes_stream_mistral(chunk, previous_summary):
                                yield response
                            previous_summary = response
                        elif model == "gemini":
                            async for response in generate_gemini_notes_stream(chunk, previous_summary):
                                yield response
                            previous_summary = response
            except Exception as e:
                logging.error(f"Streaming error: {e}")
                yield f"Error: {str(e)}"
