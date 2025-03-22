import os
import spacy
import chromadb
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_huggingface import HuggingFaceEmbeddings 
from core.config import GEMINI_API_KEY as GOOGLE_API_KEY, OPENAI_API_KEY,MISTRAL_API_KEY
from core.prompts import CHAT_PROMPT
import os
import pdfplumber
from pptx import Presentation
import pandas as pd
from mistralai import Mistral
import base64
from langchain_core.retrievers import BaseRetriever

nlp = spacy.load("en_core_web_sm")

class DocumentChatServiceGemini:
    def __init__(self, collection_name="document-chat-collection-gemini"):
        self.llm = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash",
            google_api_key=GOOGLE_API_KEY,
            temperature=0.7
        )
        self.embeddings = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"}
        )

        self.memory = ConversationBufferMemory(
            memory_key="chat_history",
            output_key="answer",
            return_messages=True
        )
        self.client = chromadb.Client()
        self.collection_name = collection_name
        self.collection = self.client.get_or_create_collection(
            name=self.collection_name,
            metadata={"hnsw:space": "cosine"}
        )

        self.custom_prompt = PromptTemplate(
            input_variables=["context", "question"],
            template=CHAT_PROMPT
        )
    def encode_image(self,image_path):
        """Encode the image to base64."""
        try:
            with open(image_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except FileNotFoundError:
            print(f"Error: The file {image_path} was not found.")
            return None
        except Exception as e:  
            print(f"Error: {e}")
            return None

    def extract_text_from_file(self,file_path: str) -> str:
        """Extract text from different file types."""
        ext = os.path.splitext(file_path)[-1].lower()
        text = ""
        
        try:
            if ext == ".pdf":
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join([page.extract_text() or "" for page in pdf.pages])
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif ext in [".jpg", ".png"]:
                mistral_client = Mistral(api_key=MISTRAL_API_KEY)
                base64_image = self.encode_image(file_path)
                ocr_response = mistral_client.ocr.process(
                    model="mistral-ocr-latest",
                    document={"type": "image_url", "image_url": f"data:image/jpeg;base64,{base64_image}"}
                )
                text = "\n".join(page.markdown for page in ocr_response.pages)
            elif ext == ".pptx":
                prs = Presentation(file_path)
                text = "\n".join([shape.text.strip() for slide in prs.slides 
                                for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
            elif ext in [".xlsx", ".csv"]:
                df = pd.read_excel(file_path, dtype=str) if ext == ".xlsx" else pd.read_csv(file_path, dtype=str)
                text = df.to_string(index=False)
            else:
                text = "Unsupported file type."
        except Exception as e:
            text = f"Error processing file: {str(e)}"
        
        return text.strip()

    def load_file(self,file_path: str):
        """Loads and processes a file of any supported type, replacing existing data in Chroma."""
        # Extract text using the existing function
        text = self.extract_text_from_file(file_path)
        
        if not text:
            print("No text extracted from the file.")
            return

        # Create a single Document object from the extracted text
        doc = Document(page_content=text, metadata={"source": file_path})
        
        # Split the text into chunks
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents([doc])
        
        if not chunks:
            print("No text chunks extracted from the file.")
            return

        self._upsert_to_chroma(chunks)

    def _upsert_to_chroma(self, docs):
        """Upserts document chunks to Chroma."""
        embeddings = [self.embeddings.embed_query(doc.page_content) for doc in docs]
        ids = [f"doc_{i}" for i in range(len(docs))]
        metadatas = [{"text": doc.page_content} for doc in docs]
        documents = [doc.page_content for doc in docs]

        self.client.delete_collection(self.collection_name)
        self.collection = self.client.create_collection(name=self.collection_name, metadata={"hnsw:space": "cosine"})
        self.collection.add(
            embeddings=embeddings,
            ids=ids,
            metadatas=metadatas,
            documents=documents
        )
        print(f"Upserted {len(embeddings)} vectors to Chroma (Gemini). Collection count: {self.collection.count()}")

    def ask_question(self, query: str):
        """Answers queries using Chroma vector search."""
        query_embedding = self.embeddings.embed_query(query)
        search_results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=10,
            include=["metadatas", "documents"]
        )
        print(f"Search results (Gemini): {search_results}")

        if not search_results["ids"][0] or not search_results["documents"][0]:
            print("No documents found in Chroma search (Gemini).")
            return {"answer": "I couldn't find relevant information in the uploaded document."}

        valid_docs = [doc for doc in search_results["documents"][0] if doc is not None]
        if not valid_docs:
            print("All documents in search results were None (Gemini).")
            return {"answer": "I couldn't find relevant information in the uploaded document."}
        context = "\n".join(valid_docs)
        print(f"Context from Chroma (Gemini): {context[:100]}...")

        class ChromaRetriever(BaseRetriever):
            def __init__(self, collection, embeddings):
                super().__init__()
                self._collection = collection
                self._embeddings = embeddings

            def _get_relevant_documents(self, query: str):
                query_embedding = self._embeddings.embed_query(query)
                results = self._collection.query(
                    query_embeddings=[query_embedding],
                    n_results=10,
                    include=["metadatas", "documents"]
                )
                valid_docs = [doc for doc in results["documents"][0] if doc is not None]
                return [Document(page_content=doc, metadata={"text": meta["text"]})
                        for doc, meta in zip(valid_docs, results["metadatas"][0])]

        retriever = ChromaRetriever(self.collection, self.embeddings)
        conversation_chain = ConversationalRetrievalChain.from_llm(
            llm=self.llm,
            retriever=retriever,
            memory=self.memory,
            return_source_documents=True,
            combine_docs_chain_kwargs={"prompt": self.custom_prompt}
        )

        response = conversation_chain.invoke({"question": query, "chat_history": self.memory.load_memory_variables({})["chat_history"]})
        return {"answer": response.get("answer", "I couldn't find relevant information in the uploaded document.")}
    
class DocumentChatServiceOpenAI:
    def __init__(self, collection_name="document-chat-collection-openai"):
        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            openai_api_key=OPENAI_API_KEY,
            temperature=0.7
        )
        self.embeddings = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"}
        )
        self.memory = ConversationBufferMemory(
            memory_key="chat_history",
            output_key="answer",
            return_messages=True
        )
        self.client = chromadb.Client()
        self.collection_name = collection_name
        self.collection = self.client.get_or_create_collection(
            name=self.collection_name,
            metadata={"hnsw:space": "cosine"}
        )

        self.custom_prompt = PromptTemplate(
            input_variables=["context", "question"],
            template=CHAT_PROMPT)

    def encode_image(self,image_path):
        """Encode the image to base64."""
        try:
            with open(image_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except FileNotFoundError:
            print(f"Error: The file {image_path} was not found.")
            return None
        except Exception as e:  
            print(f"Error: {e}")
            return None

    def extract_text_from_file(self,file_path: str) -> str:
        """Extract text from different file types."""
        ext = os.path.splitext(file_path)[-1].lower()
        text = ""
        
        try:
            if ext == ".pdf":
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join([page.extract_text() or "" for page in pdf.pages])
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif ext in [".jpg", ".png"]:
                mistral_client = Mistral(api_key=MISTRAL_API_KEY)
                base64_image = self.encode_image(file_path)
                ocr_response = mistral_client.ocr.process(
                    model="mistral-ocr-latest",
                    document={"type": "image_url", "image_url": f"data:image/jpeg;base64,{base64_image}"}
                )
                text = "\n".join(page.markdown for page in ocr_response.pages)
            elif ext == ".pptx":
                prs = Presentation(file_path)
                text = "\n".join([shape.text.strip() for slide in prs.slides 
                                for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
            elif ext in [".xlsx", ".csv"]:
                df = pd.read_excel(file_path, dtype=str) if ext == ".xlsx" else pd.read_csv(file_path, dtype=str)
                text = df.to_string(index=False)
            else:
                text = "Unsupported file type."
        except Exception as e:
            text = f"Error processing file: {str(e)}"
        
        return text.strip()

    def load_file(self,file_path: str):
        """Loads and processes a file of any supported type, replacing existing data in Chroma."""
        # Extract text using the existing function
        text = self.extract_text_from_file(file_path)
        
        if not text:
            print("No text extracted from the file.")
            return

        # Create a single Document object from the extracted text
        doc = Document(page_content=text, metadata={"source": file_path})
        
        # Split the text into chunks
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents([doc])
        
        if not chunks:
            print("No text chunks extracted from the file.")
            return

        self._upsert_to_chroma(chunks)

    def _upsert_to_chroma(self, docs):
        """Upserts document chunks to Chroma."""
        embeddings = [self.embeddings.embed_query(doc.page_content) for doc in docs]
        ids = [f"doc_{i}" for i in range(len(docs))]
        metadatas = [{"text": doc.page_content} for doc in docs]
        documents = [doc.page_content for doc in docs]

        self.client.delete_collection(self.collection_name)
        self.collection = self.client.create_collection(name=self.collection_name, metadata={"hnsw:space": "cosine"})
        self.collection.add(
            embeddings=embeddings,
            ids=ids,
            metadatas=metadatas,
            documents=documents
        )
        print(f"Upserted {len(embeddings)} vectors to Chroma (OpenAI). Collection count: {self.collection.count()}")

    def ask_question(self, query: str):
        """Answers queries using Chroma vector search."""
        query_embedding = self.embeddings.embed_query(query)
        search_results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=10,
            include=["metadatas", "documents"]
        )
        print(f"Search results (OpenAI): {search_results}")

        if not search_results["ids"][0] or not search_results["documents"][0]:
            print("No documents found in Chroma search (OpenAI).")
            return {"answer": "I couldn't find relevant information in the uploaded document."}

        valid_docs = [doc for doc in search_results["documents"][0] if doc is not None]
        if not valid_docs:
            print("All documents in search results were None (OpenAI).")
            return {"answer": "I couldn't find relevant information in the uploaded document."}
        context = "\n".join(valid_docs)
        print(f"Context from Chroma (OpenAI): {context[:100]}...")

        from langchain_core.retrievers import BaseRetriever
        class ChromaRetriever(BaseRetriever):
            def __init__(self, collection, embeddings):
                super().__init__()
                self._collection = collection
                self._embeddings = embeddings

            def _get_relevant_documents(self, query: str):
                query_embedding = self._embeddings.embed_query(query)
                results = self._collection.query(
                    query_embeddings=[query_embedding],
                    n_results=10,
                    include=["metadatas", "documents"]
                )
                valid_docs = [doc for doc in results["documents"][0] if doc is not None]
                return [Document(page_content=doc, metadata={"text": meta["text"]})
                        for doc, meta in zip(valid_docs, results["metadatas"][0])]

        retriever = ChromaRetriever(self.collection, self.embeddings)
        conversation_chain = ConversationalRetrievalChain.from_llm(
            llm=self.llm,
            retriever=retriever,
            memory=self.memory,
            return_source_documents=True,
            combine_docs_chain_kwargs={"prompt": self.custom_prompt}
        )

        response = conversation_chain.invoke({"question": query, "chat_history": self.memory.load_memory_variables({})["chat_history"]})
        return {"answer": response.get("answer", "I couldn't find relevant information in the uploaded document.")}
